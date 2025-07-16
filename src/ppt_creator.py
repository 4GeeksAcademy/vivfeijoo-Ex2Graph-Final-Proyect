from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import os

# This function saves the current matplotlib plot and inserts it into PowerPoint
def save_plot_to_ppt(image_path, ppt_path):
    # Create a PowerPoint presentation
    prs = Presentation()

    # Add a blank slide
    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add the image to the slide
    slide.shapes.add_picture(image_path, Inches(1), Inches(1), width=Inches(8))

    # Save the PowerPoint
    prs.save(ppt_path)
    print(f"✅ PowerPoint saved to: {ppt_path}")


# This function creates a simple example chart and saves + inserts it
def create_chart_and_ppt():
    # Example data
    labels = ["On Time", "Late"]
    values = [8, 2]

    # Create a basic chart
    plt.figure()
    plt.bar(labels, values, color=["green", "red"])
    plt.title("Delivery Status")
    plt.tight_layout()

    # Paths
    image_path = "output/temp_chart.png"
    ppt_path = "output/delivery_report.pptx"

    # Save chart
    os.makedirs("output", exist_ok=True)
    plt.savefig(image_path)
    plt.close()

    # Insert into PowerPoint
    save_plot_to_ppt(image_path, ppt_path)
